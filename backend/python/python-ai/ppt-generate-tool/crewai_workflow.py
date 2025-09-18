import os
from typing import List, Dict, Any, Union
from crewai import Agent, Task, Crew, Process, LLM
from ddgs import DDGS
from pptx import Presentation
from pptx.util import Pt
from config import Config

# -------------------------
# Helper: robust crew output extractor
# -------------------------
def extract_text(crew_output: Any) -> str:
    """Try several fields to get a readable string from a Crew kickoff result."""
    if crew_output is None:
        return ""
    if isinstance(crew_output, str):
        return crew_output
    # try common attributes
    for attr in ("final_output", "raw", "output", "result"):
        val = getattr(crew_output, attr, None)
        if isinstance(val, str) and val.strip():
            return val.strip()
    # tasks_output can be a list
    tasks = getattr(crew_output, "tasks_output", None) or getattr(crew_output, "outputs", None)
    if tasks:
        parts = []
        if isinstance(tasks, (list, tuple)):
            for t in tasks:
                if isinstance(t, str):
                    parts.append(t.strip())
                elif isinstance(t, dict):
                    for k in ("final_output", "output", "result", "text"):
                        if k in t and isinstance(t[k], str):
                            parts.append(t[k].strip())
                            break
        if parts:
            return "\n\n".join(parts)
    return str(crew_output)


# -------------------------
# LLM setup (CrewAI LLM wrapper)
# -------------------------
def get_llm():
    # Use GEMINI_API_KEY if present (CrewAI will pass through to provider)
    if getattr(Config, "GEMINI_API_KEY", None):
        return LLM(
            model="gemini/gemini-1.5-flash",
            api_key=Config.GEMINI_API_KEY,
            temperature=0.2,
        )
    else:
        # Dummy fallback (simple echo)
        class DummyLLM:
            def generate(self, prompt: str) -> str:
                return f"[Dummy LLM] {prompt}"
        return DummyLLM()

LLM_CLIENT = get_llm()


# -------------------------
# Agents
# -------------------------
chat_agent = Agent(
    role="Chat Assistant",
    goal="Handle general conversational queries and decide whether an external search is required.",
    backstory="You are the user-facing assistant. If research is necessary respond with 'search: <query>' in your output.",
    llm=LLM_CLIENT,
    verbose=True,
)

summarizer_agent = Agent(
    role="Summarizer",
    goal="Create a concise, structured summary from provided search results.",
    backstory="You turn multiple search results into a short, clear summary the user can review.",
    llm=LLM_CLIENT,
    verbose=True,
)

# Note: we rely on direct functions for search and PPT generation to avoid tool API mismatches


# -------------------------
# Search helpers (ddgs)
# -------------------------
def ddg_search_raw(query: str, max_results: int = 5, timeout: int = 10) -> List[Dict[str, str]]:
    """Return list of dicts: {title, body, href}"""
    try:
        ddgs = DDGS()
        items = list(ddgs.text(query=query, max_results=max_results, timeout=timeout))
    except Exception as e:
        return [{"title": "Search error", "body": str(e), "href": ""}]
    results = []
    for it in items:
        results.append({
            "title": it.get("title", "No Title"),
            "body": it.get("body", "No snippet"),
            "href": it.get("href", "")
        })
    return results

def format_results_for_chat(results: List[Dict[str, str]]) -> str:
    if not results:
        return "No results found."
    out = []
    for i, r in enumerate(results, 1):
        out.append(f"{i}. {r['title']}\n   {r['body']}\n   {r['href']}")
    return "\n\n".join(out)


# -------------------------
# Workflow steps (Crew-driven for LLM tasks)
# -------------------------
def llm_chat_step(user_msg: str) -> str:
    """LLM handles general chat. If it thinks research is required include 'search: <query>' in reply."""
    task = Task(
        description=f"Respond conversationally to the user message: {user_msg}",
        expected_output="A conversational reply. If research is needed include 'search: <query>'.",
        agent=chat_agent,
    )
    crew = Crew(agents=[chat_agent], tasks=[task], process=Process.sequential)
    res = crew.kickoff()
    return extract_text(res)


def first_search_step(query: str) -> Dict[str, Any]:
    """Perform search (ddgs) and return both raw results and a user-facing message."""
    raw = ddg_search_raw(query, max_results=getattr(Config, "SEARCH_MAX_RESULTS", 5),
                         timeout=getattr(Config, "SEARCH_TIMEOUT_SEC", 10))
    msg = f"Top {len(raw)} results for: {query}\n\n{format_results_for_chat(raw)}\n\nReply with 'OK' to confirm these results or type a new query to search again."
    return {"results": raw, "message": msg}


def confirm_and_summarize_step(query: str, results: List[Dict[str, str]]) -> Dict[str, str]:
    """Use CrewAI summarizer agent to produce a succinct summary of given results."""
    # Build a compact text block as input to LLM
    results_text = "\n\n".join(f"{r['title']}: {r['body']} ({r['href']})" for r in results)
    prompt_desc = f"Summarize the following search results for '{query}' into a concise summary (bullet points or short paragraphs):\n\n{results_text}"
    task = Task(
        description=prompt_desc,
        expected_output="A concise summary of the search results.",
        agent=summarizer_agent,
    )
    crew = Crew(agents=[summarizer_agent], tasks=[task], process=Process.sequential)
    res = crew.kickoff()
    summary = extract_text(res)
    message = f"Summary for: {query}\n\n{summary}\n\nIf you'd like a PPT, say 'Generate PPT'."
    return {"summary": summary, "message": message}


def generate_ppt_step(summary: str) -> Dict[str, str]:
    """Create a PPTX from the summary and return the file path."""
    os.makedirs("outputs", exist_ok=True)
    output_file = os.path.join("outputs", "presentation.pptx")

    prs = Presentation()
    # Title slide
    try:
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Research Summary"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Generated by CrewAI"
    except Exception:
        # fallback: create a basic slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Content slide
    try:
        bullet_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Key Points"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for line in summary.splitlines():
            t = line.strip("- •\t ").strip()
            if not t:
                continue
            p = tf.add_paragraph()
            p.text = t
            p.level = 0
            p.font.size = Pt(18)
    except Exception:
        # fallback: append summary in a blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # attempt to put text in title shape if present
        if slide.shapes.title:
            slide.shapes.title.text = "Summary"
            # no detailed formatting in fallback

    prs.save(output_file)
    return {"ppt_path": output_file}
